"""
Build a personal research deck from an Instagram-focused profile.md.

Usage:
    python build_instagram_deck.py <client_slug>

Example:
    python build_instagram_deck.py aidan-pinter
"""

import sys
import re
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

CREAM = RGBColor(0xF3, 0xF1, 0xE7)
TERRACOTTA = RGBColor(0xD9, 0x77, 0x57)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
LIGHT_TEXT = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _parse_markdown_runs(text):
    segments = []
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            segments.append((part[2:-2], True))
        elif part:
            segments.append((part, False))
    return segments


def _set_runs(p, text, font_size, color, font_name="Calibri"):
    segments = _parse_markdown_runs(text)
    if len(segments) <= 1 and not any(b for _, b in segments):
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        return
    p.text = ""
    for seg_text, is_bold in segments:
        run = p.add_run()
        run.text = seg_text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_name
        run.font.bold = is_bold


def add_text(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if bold:
        p.text = text
        p.font.bold = True
    else:
        _set_runs(p, text, font_size, color)
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox


def add_bullets(slide, left, top, width, height, items, font_size=14, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_runs(p, item, font_size, color)
        p.space_after = Pt(6)
    return txBox


def bar(slide, top=Inches(0), height=Inches(0.08)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_WIDTH, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TERRACOTTA
    shape.line.fill.background()


def header(slide, title, bg=WHITE):
    set_slide_bg(slide, bg)
    bar(slide)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             title, font_size=32, bold=True)


def add_picture_scaled(slide, img_path, left, top, max_w, max_h):
    if not img_path.exists():
        return
    with Image.open(str(img_path)) as img:
        iw, ih = img.size
    aspect = iw / ih
    tw = max_w
    th = int(tw / aspect)
    if th > max_h:
        th = max_h
        tw = int(th * aspect)
    slide.shapes.add_picture(str(img_path), left, top, width=tw, height=th)


def extract_section_items(text, heading):
    items = []
    in_section = False
    for line in text.split("\n"):
        if line.strip().startswith(heading):
            in_section = True
            continue
        if in_section:
            if line.startswith("## ") and heading not in line:
                break
            stripped = line.strip()
            if stripped.startswith("- "):
                items.append(stripped[2:])
            elif re.match(r"\d+\.\s+", stripped):
                items.append(re.sub(r"^\d+\.\s+", "", stripped))
    return items


def extract_table_rows(text, heading):
    rows = []
    in_section = False
    past_header = False
    for line in text.split("\n"):
        if line.strip().startswith(heading):
            in_section = True
            continue
        if in_section:
            if line.startswith("## ") and heading not in line:
                break
            stripped = line.strip()
            if stripped.startswith("|") and "---" in stripped:
                past_header = True
                continue
            if past_header and stripped.startswith("|"):
                cells = [c.strip() for c in stripped.split("|")[1:-1]]
                rows.append(cells)
    return rows


def build_deck(slug):
    base = Path(__file__).resolve().parent.parent.parent.parent / "context" / "clients" / slug
    profile_path = base / "profile.md"

    if not profile_path.exists():
        print(f"Error: {profile_path} not found")
        sys.exit(1)

    text = profile_path.read_text(encoding="utf-8")

    # Parse basics
    name_m = re.search(r"# Client Profile:\s*(.+)", text)
    name = name_m.group(1).strip() if name_m else slug.replace("-", " ").title()

    date_m = re.search(r"\*\*Researched\*\*:\s*(.+)", text)
    date = date_m.group(1).strip() if date_m else "2026-04-05"

    handle_m = re.search(r"\*\*Instagram\*\*:\s*(.+)", text)
    ig_url = handle_m.group(1).strip() if handle_m else ""

    # Extract sections
    takeaways = extract_section_items(text, "## Key Takeaways")
    posts = extract_table_rows(text, "## Posts Reviewed")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ========== Slide 1: Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)
    bar(slide, top=Inches(0), height=Inches(0.12))
    add_text(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
             name, font_size=44, bold=True)
    add_text(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
             "Deep Personalization Research", font_size=28, color=TERRACOTTA)
    add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             f"{date}  |  Instagram: @{slug}", font_size=16, color=LIGHT_TEXT)
    bar(slide, top=Inches(7.38), height=Inches(0.12))

    # ========== Slide 2: Profile Overview + Screenshot ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Instagram Profile")

    # Extract IG profile table
    ig_rows = extract_table_rows(text, "## Instagram Profile")
    overview = []
    for row in ig_rows:
        if len(row) >= 2:
            overview.append(f"**{row[0]}**: {row[1]}")
    if not overview:
        overview = [f"**Handle**: @{slug}", f"**Name**: {name}"]

    add_bullets(slide, Inches(0.8), Inches(1.4), Inches(6.5), Inches(5.5),
                overview, font_size=16)

    # Profile screenshot on right
    add_picture_scaled(slide, base / "instagram-profile.png",
                       Inches(7.8), Inches(1.2), Inches(5), Inches(5.5))

    # ========== Slide 3: Engagement & Post Analysis ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Engagement Analysis")

    # Extract engagement section
    eng_items = extract_section_items(text, "## Engagement Analysis")
    if eng_items:
        add_bullets(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(3),
                    eng_items, font_size=15)

    # Recent posts summary on right
    add_text(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
             "Recent Posts", font_size=20, bold=True, color=TERRACOTTA)
    post_lines = []
    for row in posts[:6]:
        if len(row) >= 4:
            date_str = row[1]
            caption = row[2].strip('"')
            if len(caption) > 50:
                caption = caption[:47] + "..."
            likes = row[3]
            post_lines.append(f"**{date_str}** — \"{caption}\" ({likes} likes)")
    if post_lines:
        add_bullets(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(5),
                    post_lines, font_size=12)

    # ========== Slide 4: Caption Intelligence ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Caption Intelligence", bg=CREAM)

    # Left: Tone & Travel
    add_text(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
             "Tone & Style", font_size=20, bold=True, color=TERRACOTTA)
    tone_items = extract_section_items(text, "### Tone & Style")
    if tone_items:
        add_bullets(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
                    tone_items[:4], font_size=13)

    add_text(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.5),
             "Travel & Adventure", font_size=20, bold=True, color=TERRACOTTA)
    travel_items = extract_section_items(text, "### Travel & Adventure")
    if travel_items:
        add_bullets(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2),
                    travel_items[:4], font_size=13)

    # Right: Family & Inner Circle
    add_text(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
             "Family & Relationships", font_size=20, bold=True, color=TERRACOTTA)
    family_items = extract_section_items(text, "### Family & Relationships")
    if family_items:
        add_bullets(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
                    family_items[:4], font_size=13)

    add_text(slide, Inches(7), Inches(4.2), Inches(5.5), Inches(0.5),
             "Inner Circle", font_size=20, bold=True, color=TERRACOTTA)
    circle_items = extract_section_items(text, "### Inner Circle")
    if circle_items:
        add_bullets(slide, Inches(7), Inches(4.8), Inches(5.5), Inches(2),
                    circle_items[:4], font_size=13)

    # ========== Slide 5: Visual Grid (post screenshots) ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Post Highlights")

    # 2x3 grid of post screenshots
    post_imgs = sorted(base.glob("instagram-post-*.png"))[:6]
    positions = [
        (Inches(0.5), Inches(1.3)),  (Inches(4.5), Inches(1.3)),  (Inches(8.5), Inches(1.3)),
        (Inches(0.5), Inches(4.3)),  (Inches(4.5), Inches(4.3)),  (Inches(8.5), Inches(4.3)),
    ]
    for img_path, (left, top) in zip(post_imgs, positions):
        add_picture_scaled(slide, img_path, left, top, Inches(3.8), Inches(2.8))

    # ========== Slide 6: Key Takeaways ==========
    if takeaways:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        header(slide, "Key Takeaways", bg=CREAM)
        add_bullets(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                    takeaways[:8], font_size=15)

    # ========== Slide 7: Next Steps — Conversation Starters ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)
    bar(slide)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Next Steps & Conversation Starters", font_size=32, bold=True)

    add_text(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
             "Conversation Openers", font_size=20, bold=True, color=TERRACOTTA)

    # Build starters from the data
    starters = []
    values_items = extract_section_items(text, "### Values Signaled")
    if any("travel" in t.lower() for t in takeaways + values_items):
        starters.append("\"31 days in the air\" — where was the best trip this year?")
    if any("diplo" in t.lower() for t in takeaways):
        starters.append("The Diplo connection — how do you know each other?")
    if any("family" in t.lower() or "wedding" in t.lower() for t in takeaways):
        starters.append("Dad's wedding looked incredible — how's the blended family?")
    if any("one-liner" in t.lower() or "witty" in t.lower() or "humor" in t.lower() for t in takeaways):
        starters.append("Your captions are legendary — \"afraid of farm animals\" is gold")
    if any("engagement" in t.lower() for t in takeaways):
        starters.append("17.7% engagement rate — your audience genuinely loves your content")
    starters.append("What's the story behind no bio? Intentional minimalism?")
    starters.append("The motorcycle photo — Vietnam? Southeast Asia?")

    add_bullets(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
                starters[:6], font_size=14)

    add_text(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
             "Creative Personalization Angles", font_size=20, bold=True, color=TERRACOTTA)

    angles = [
        "**Adventure-first identity** — lead with experiences, not credentials",
        "**Humor as currency** — match his witty, minimal tone in outreach",
        "**Celebrity adjacency** — Diplo comments = real social capital",
        "**Family values** — emotionally intelligent, values loyalty",
        "**Curated minimalism** — quality over quantity in everything",
        "**No bio = no labels** — he defines himself through actions, not titles",
    ]
    add_bullets(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(5),
                angles, font_size=14)

    bar(slide, top=Inches(7.38), height=Inches(0.12))

    # Save
    out_path = base / f"{name} - Personal Intel.pptx"
    try:
        prs.save(str(out_path))
    except PermissionError:
        from datetime import datetime
        ts = datetime.now().strftime("%H%M%S")
        out_path = base / f"{name} - Personal Intel {ts}.pptx"
        prs.save(str(out_path))
    print(f"Deck saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    return str(out_path)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python build_instagram_deck.py <client-slug>")
        sys.exit(1)
    build_deck(sys.argv[1])
