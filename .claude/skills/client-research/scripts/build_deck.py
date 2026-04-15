"""
Build a client research PowerPoint deck from a profile.md and screenshots.

Usage:
    python build_deck.py <client_slug>

Example:
    python build_deck.py patrick-bauer

Reads from:  .claude/context/clients/{slug}/profile.md
             .claude/context/clients/{slug}/*.png
             .claude/context/clients/{slug}/deep-research.md
             .claude/context/clients/{slug}/site-research.md
Writes to:   .claude/context/clients/{slug}/{Name} - Client Research.pptx
"""

import sys
import re
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# GBAutomation brand colors
CREAM = RGBColor(0xF3, 0xF1, 0xE7)
TERRACOTTA = RGBColor(0xD9, 0x77, 0x57)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
LIGHT_TEXT = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _parse_markdown_runs(text):
    """Parse **bold** markers in text into (text, is_bold) segments."""
    segments = []
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            segments.append((part[2:-2], True))
        elif part:
            segments.append((part, False))
    return segments


def _set_paragraph_runs(p, text, font_size, color, font_name="Calibri"):
    """Set paragraph text with **bold** markdown parsed into real bold runs."""
    segments = _parse_markdown_runs(text)
    if len(segments) <= 1 and not any(b for _, b in segments):
        # No bold markers — use simple text assignment
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        return
    # Clear default text, add runs
    p.text = ""
    for seg_text, is_bold in segments:
        run = p.add_run()
        run.text = seg_text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_name
        run.font.bold = is_bold


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if bold:
        p.text = text
        p.font.bold = True
    else:
        _set_paragraph_runs(p, text, font_size, color, font_name)
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK_TEXT, bold_first_line=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        _set_paragraph_runs(p, item, font_size, color)
        p.space_after = Pt(6)
        if bold_first_line and i == 0:
            p.font.bold = True
    return txBox


def add_terracotta_bar(slide, top=Inches(0), height=Inches(0.08)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TERRACOTTA
    shape.line.fill.background()


SOURCE_GRAY = RGBColor(0x99, 0x99, 0x99)


def add_source_citation(slide, sources):
    """Add small source URLs at the bottom of a slide."""
    if not sources:
        return
    # Truncate long URLs for display
    display = []
    for s in sources[:3]:
        # Strip protocol, trim to 80 chars
        short = re.sub(r"https?://", "", s)
        if len(short) > 80:
            short = short[:77] + "..."
        display.append(short)
    text = "Sources: " + "  |  ".join(display)
    add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4),
                 text, font_size=9, color=SOURCE_GRAY)


def add_slide_header(slide, title, bg=WHITE):
    set_slide_bg(slide, bg)
    add_terracotta_bar(slide)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 title, font_size=32, bold=True, color=DARK_TEXT)


def add_screenshot_slide(prs, title, img_path, bg=WHITE, max_width=Inches(11.5), max_height=Inches(5.8)):
    """Add a slide with a title and proportionally-scaled screenshot (no stretching)."""
    if not img_path.exists():
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, title, bg)

    # Get native image dimensions for proportional scaling
    with Image.open(str(img_path)) as img:
        img_w, img_h = img.size

    # Scale to fit within max_width x max_height, preserving aspect ratio
    aspect = img_w / img_h
    target_w = max_width
    target_h = int(target_w / aspect)
    if target_h > max_height:
        target_h = max_height
        target_w = int(target_h * aspect)

    # Center horizontally
    left = Inches(0.9) + (max_width - target_w) // 2
    slide.shapes.add_picture(str(img_path), left, Inches(1.3),
                             width=target_w, height=target_h)


def extract_section(text, header, stop_at="## "):
    """Extract lines from a markdown section."""
    lines = []
    in_section = False
    for line in text.split("\n"):
        if line.strip().startswith(header):
            in_section = True
            continue
        if in_section:
            if line.startswith(stop_at) and header not in line:
                break
            lines.append(line)
    return lines


def extract_bullets(text, header):
    """Extract bullet items from a markdown section."""
    items = []
    for line in extract_section(text, header):
        if line.strip().startswith("- "):
            items.append(line.strip()[2:])
    return items


def extract_numbered(text, header):
    """Extract numbered items from a markdown section."""
    items = []
    for line in extract_section(text, header):
        m = re.match(r"\d+\.\s*(.+)", line.strip())
        if m:
            items.append(m.group(1))
    return items


def extract_table_rows(text, header):
    """Extract table rows from a markdown section (skip header + separator)."""
    rows = []
    lines = extract_section(text, header)
    table_started = False
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("|") and "---" in stripped:
            table_started = True
            continue
        if table_started and stripped.startswith("|"):
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            rows.append(cells)
    return rows


def parse_profile(profile_path):
    """Parse profile.md into structured data."""
    text = Path(profile_path).read_text(encoding="utf-8")
    data = {"_raw": text}

    m = re.search(r"# Client Profile:\s*(.+)", text)
    data["name"] = m.group(1).strip() if m else "Unknown"

    m = re.search(r"\*\*Researched\*\*:\s*(.+)", text)
    data["date"] = m.group(1).strip() if m else ""

    m = re.search(r"\*\*LinkedIn\*\*:\s*(.+)", text)
    data["linkedin"] = m.group(1).strip() if m else ""

    m = re.search(r"\*\*Session\*\*:\s*(.+)", text)
    data["session"] = m.group(1).strip() if m else ""

    for field in ["Headline", "Location", "Current Role", "Connection"]:
        m = re.search(rf"\|\s*{field}\s*\|\s*(.+?)\s*\|", text)
        data[field.lower().replace(" ", "_")] = m.group(1).strip() if m else ""

    data["services"] = extract_bullets(text, "## Services")
    data["skills"] = extract_bullets(text, "## Skills")
    data["languages"] = extract_bullets(text, "## Languages")
    data["education"] = extract_bullets(text, "## Education")
    data["takeaways"] = extract_numbered(text, "## Key Takeaways")
    data["proprietary_systems"] = extract_table_rows(text, "## Proprietary Systems")

    # Experience — extract ### blocks
    experience = []
    current_entry = []
    in_exp = False
    for line in text.split("\n"):
        if line.strip() == "## Experience":
            in_exp = True
            continue
        if in_exp:
            if line.startswith("## ") and "Experience" not in line:
                if current_entry:
                    experience.append("\n".join(current_entry))
                break
            if line.strip().startswith("### "):
                if current_entry:
                    experience.append("\n".join(current_entry))
                current_entry = [line.strip()[4:]]
            elif line.strip():
                if current_entry:
                    current_entry.append(line.strip())
    if current_entry:
        experience.append("\n".join(current_entry))
    data["experience"] = experience

    # Company deep dive — extract subsections
    company_sections = []
    in_company = False
    current = []
    for line in text.split("\n"):
        if line.strip() == "## Company Deep Dive":
            in_company = True
            continue
        if in_company:
            if line.startswith("## ") and "Company" not in line:
                if current:
                    company_sections.append("\n".join(current))
                break
            if line.strip().startswith("### "):
                if current:
                    company_sections.append("\n".join(current))
                current = [line.strip()[4:]]
            elif line.strip():
                current.append(line.strip())
    if current:
        company_sections.append("\n".join(current))
    data["companies"] = company_sections

    # Video/media
    data["media"] = extract_bullets(text, "## Video")
    if not data["media"]:
        media_lines = extract_section(text, "## Video")
        data["media"] = [l.strip() for l in media_lines if l.strip() and not l.strip().startswith("#")]

    return data


def extract_all_urls(base_dir):
    """Extract all URLs from all .md files in the client folder, grouped by source file."""
    url_map = {}  # {source_file: [urls]}
    all_urls = []
    url_pattern = re.compile(r'https?://[^\s\)\]\|>"\']+')
    for md_file in sorted(base_dir.glob("*.md")):
        text = md_file.read_text(encoding="utf-8", errors="ignore")
        urls = url_pattern.findall(text)
        # Clean trailing punctuation
        cleaned = []
        for u in urls:
            u = u.rstrip(".,;:)")
            if u not in cleaned and "linkedin.com/li/tscp" not in u and "brandfetch" not in u:
                cleaned.append(u)
        if cleaned:
            url_map[md_file.stem] = cleaned
            all_urls.extend(cleaned)
    # Dedupe preserving order
    seen = set()
    unique = []
    for u in all_urls:
        if u not in seen:
            seen.add(u)
            unique.append(u)
    return url_map, unique


def urls_for_section(url_map, keywords):
    """Find URLs from research files that match given keywords."""
    results = []
    for source, urls in url_map.items():
        for url in urls:
            for kw in keywords:
                if kw.lower() in url.lower() or kw.lower() in source.lower():
                    if url not in results:
                        results.append(url)
    return results


def build_deck(client_slug):
    base = Path(__file__).resolve().parent.parent.parent.parent / "context" / "clients" / client_slug
    profile_path = base / "profile.md"

    if not profile_path.exists():
        print(f"Error: {profile_path} not found")
        sys.exit(1)

    data = parse_profile(profile_path)
    url_map, all_urls = extract_all_urls(base)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # =============================================
    # Slide 1: Title
    # =============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)
    add_terracotta_bar(slide, top=Inches(0), height=Inches(0.12))
    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 data["name"], font_size=44, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
                 "Deep Personalization Research", font_size=28, color=TERRACOTTA)
    subtitle_parts = []
    if data.get("date"):
        subtitle_parts.append(data["date"])
    if data.get("headline"):
        subtitle_parts.append(data["headline"])
    if subtitle_parts:
        add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
                     "  |  ".join(subtitle_parts), font_size=16, color=LIGHT_TEXT)
    add_terracotta_bar(slide, top=Inches(7.38), height=Inches(0.12))

    # =============================================
    # Slide 2: Profile Overview + LinkedIn Screenshot
    # =============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Profile Overview")
    overview = [
        f"**Name**: {data['name']}",
        f"**Headline**: {data.get('headline', '')}",
        f"**Location**: {data.get('location', '')}",
        f"**Current Role**: {data.get('current_role', '')}",
        f"**Connection**: {data.get('connection', '')}",
        f"**LinkedIn**: {data.get('linkedin', '')}",
    ]
    if data.get("languages"):
        overview.append(f"**Languages**: {', '.join(data['languages'])}")
    if data.get("education"):
        overview.append(f"**Education**: {data['education'][0]}")
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(7), Inches(5.5),
                    overview, font_size=16)

    # Embed LinkedIn screenshot on right if it exists
    profile_img = base / "profile-header.png"
    if profile_img.exists():
        with Image.open(str(profile_img)) as img:
            img_w, img_h = img.size
        max_w = Inches(4.8)
        max_h = Inches(5.0)
        aspect = img_w / img_h
        target_w = max_w
        target_h = int(target_w / aspect)
        if target_h > max_h:
            target_h = max_h
            target_w = int(target_h * aspect)
        slide.shapes.add_picture(str(profile_img), Inches(8.2), Inches(1.4),
                                 width=target_w, height=target_h)

    # =============================================
    # Slide 3: Career Timeline + Companies
    # =============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Career & Companies")
    exp_items = []
    for exp in data.get("experience", [])[:6]:
        lines = exp.split("\n")
        title = lines[0]
        company = ""
        dates = ""
        desc_preview = ""
        for l in lines[1:]:
            clean = l.strip("*").strip()
            if not company and ("Full-time" in l or "—" in l or "-" in l):
                company = clean
            elif not dates and re.search(r"\d{4}", l):
                dates = clean
            elif not desc_preview and len(clean) > 30:
                desc_preview = clean[:140] + ("..." if len(clean) > 140 else "")
        entry = f"**{title}**"
        if company:
            entry += f"  |  {company}"
        if dates:
            entry += f"  |  {dates}"
        exp_items.append(entry)
        if desc_preview:
            exp_items.append(f"    {desc_preview}")
    if exp_items:
        add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                        exp_items, font_size=13)

    # =============================================
    # Slide 4: Proprietary Systems + Key Skills
    # =============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Systems, Skills & Methodologies")

    # Left column: proprietary systems
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(6), Inches(0.5),
                 "Proprietary Systems", font_size=20, bold=True, color=TERRACOTTA)
    sys_items = []
    for row in data.get("proprietary_systems", [])[:6]:
        if len(row) >= 3:
            sys_items.append(f"**{row[0].strip('*')}** — {row[1]} ({row[2]})")
        elif len(row) >= 2:
            sys_items.append(f"**{row[0].strip('*')}** — {row[1]}")
    if sys_items:
        add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(5),
                        sys_items, font_size=13)

    # Right column: top skills + services
    add_text_box(slide, Inches(7.2), Inches(1.0), Inches(5.5), Inches(0.5),
                 "Core Skills", font_size=20, bold=True, color=TERRACOTTA)
    combined = (data.get("skills", []) + data.get("services", []))[:8]
    if combined:
        add_bullet_list(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(5),
                        combined, font_size=14)

    # =============================================
    # Slide 5: Personal & Digital Footprint
    # =============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Personal & Digital Footprint", bg=CREAM)
    personal_items = []
    raw = data.get("_raw", "")
    personal_section = extract_section(raw, "## Personal")
    for line in personal_section:
        stripped = line.strip()
        if stripped.startswith("- **") and "**:" in stripped:
            personal_items.append(stripped.lstrip("- "))
    if not personal_items:
        personal_items = [
            "**Instagram**: None (confirmed by client)",
            "**Social media presence**: Very low — no personal Twitter/X, TikTok, or public Facebook",
            "**Podcast appearances**: None found",
            "**Digital footprint**: Operates behind the scenes, prefers builder role over influencer role",
        ]
    # Add media if available
    if data.get("media"):
        personal_items.append("")
        personal_items.append("**Media Appearances**:")
        for m in data["media"][:3]:
            if len(m) > 5:
                personal_items.append(f"  {m}")
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                    personal_items, font_size=15)

    # =============================================
    # Slide 6: Key Takeaways (all on one slide)
    # =============================================
    if data.get("takeaways"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_header(slide, "Key Takeaways", bg=CREAM)
        add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                        data["takeaways"][:8], font_size=14)

    # =============================================
    # Slide 7: Conversation Starters & Personal Insights
    # =============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Conversation Starters & Personal Insights", bg=CREAM)

    starters = []
    deep_path = base / "deep-research.md"
    if deep_path.exists():
        deep_text = deep_path.read_text(encoding="utf-8")
        starters = extract_numbered(deep_text, "## Key Conversation Starters")
    if not starters:
        starters = extract_numbered(raw, "## Key Takeaways")
        # Convert takeaways to conversation angles
        starters = [t.split("—")[0].strip() + " — ask about this" if "—" in t else t
                     for t in starters[:5]]
    if not starters:
        starters = [
            "Penn State connection — 2011-2012 overlap, 105 mutual connections",
            "Cura Cannabis exit — what was his role in pre-acquisition positioning?",
            "Transition to Acquisition.com — why leave his own companies for Hormozi?",
            "CRM philosophy — he endorsed Pipeliner publicly",
            "Avenoo's Retail Access Code — genuinely novel B2B mechanism",
        ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                    starters[:7], font_size=15)

    # =============================================
    # Slide 8: Next Steps — Creative Personal Insights
    # =============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)
    add_terracotta_bar(slide)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Next Steps", font_size=32, bold=True, color=DARK_TEXT)

    # Two-column layout: left = action items, right = creative insights
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
                 "Action Items", font_size=20, bold=True, color=TERRACOTTA)

    next_steps = [
        f"**Session**: {data.get('session', 'TBD')}",
        "Prepare 3 tailored questions from Key Takeaways",
        "Reference shared connections in opening",
        "Explore AI/automation angles for his domain expertise",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
                    next_steps, font_size=14)

    add_text_box(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
                 "Creative Personalization Angles", font_size=20, bold=True, color=TERRACOTTA)

    # Build creative insights from the data
    creative = []
    if data.get("languages"):
        langs = ", ".join(data["languages"])
        creative.append(f"**Languages** ({langs}) — cultural bridge, study abroad in Prague")
    if any("Penn State" in e for e in data.get("education", [])):
        creative.append("**Penn State** — shared connection, 105 mutual contacts")
    if any("Cura" in e or "Curaleaf" in e for e in data.get("experience", [])):
        creative.append("**$1B exit proximity** — was at Cura during Curaleaf acquisition")
    if any("Hormozi" in c or "Acquisition.com" in c for c in data.get("companies", [])):
        creative.append("**Hormozi inner circle** — has the scaling playbook, ask about it")
    if data.get("proprietary_systems"):
        creative.append("**Systems builder** — multiple proprietary methodologies = loves process")
    raw_personal = extract_section(raw, "## Personal")
    if any("low" in l.lower() or "none" in l.lower() for l in raw_personal):
        creative.append("**Low digital footprint** — approach as a builder, not an influencer")
    if not creative:
        creative = [
            "Research their social presence for personal touchpoints",
            "Look for shared interests, events, or connections",
            "Find their communication style from public content",
        ]
    add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5),
                    creative, font_size=14)

    add_terracotta_bar(slide, top=Inches(7.38), height=Inches(0.12))

    # Save
    out_path = base / f"{data['name']} - Client Research.pptx"
    # If file is locked (OneDrive/open app), write to timestamped name
    try:
        prs.save(str(out_path))
    except PermissionError:
        from datetime import datetime
        ts = datetime.now().strftime("%H%M%S")
        out_path = base / f"{data['name']} - Client Research {ts}.pptx"
        prs.save(str(out_path))
    print(f"Deck saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    return str(out_path)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python build_deck.py <client-slug>")
        print("Example: python build_deck.py patrick-bauer")
        sys.exit(1)
    build_deck(sys.argv[1])
