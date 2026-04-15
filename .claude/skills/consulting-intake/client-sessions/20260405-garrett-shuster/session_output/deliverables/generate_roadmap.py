"""
Generate the Sylvan Hills x GBAutomation Engagement Roadmap deck.

Usage:
    python generate_roadmap.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Brand Colors ────────────────────────────────────────────────────────────
CREAM      = RGBColor(0xF3, 0xF1, 0xE7)
TERRACOTTA = RGBColor(0xD9, 0x77, 0x57)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)
MID        = RGBColor(0x55, 0x55, 0x55)
LIGHT_TAN  = RGBColor(0xE8, 0xE5, 0xD8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TERRA_DARK = RGBColor(0xB5, 0x5E, 0x3A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

SKILL_ROOT = Path(__file__).resolve().parent.parent.parent.parent.parent / "consulting-admin"
LOGO_PATH = SKILL_ROOT / "assets" / "gb-logo.png"
OUTPUT = Path(__file__).resolve().parent / "sylvan-hills-roadmap.pptx"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_logo(slide):
    if LOGO_PATH.exists():
        from PIL import Image
        with Image.open(LOGO_PATH) as img:
            w, h = img.size
            aspect = w / h
        height = Inches(0.55)
        width = height * aspect
        left = SLIDE_W - width - Inches(0.3)
        slide.shapes.add_picture(str(LOGO_PATH), left, Inches(0.15), width, height)


def section_bar(slide, label):
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill_color=TERRACOTTA)
    add_text(slide, label, Inches(0.28), Inches(0.18), Inches(4), Inches(0.45),
             size=11, bold=True, color=TERRACOTTA)


def divider(slide, top):
    add_rect(slide, Inches(0.4), top, Inches(12.5), Pt(1.5), fill_color=LIGHT_TAN)


def bullets(slide, items, l, t, w, h, size=14):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        run = p.add_run()
        prefix = "• " if level == 0 else "   ◦ "
        run.text = prefix + text
        run.font.size = Pt(size - level * 1.5)
        run.font.color.rgb = DARK if level == 0 else MID
        run.font.name = "Calibri"
        p.space_after = Pt(4)
    return txb


# ─── Slides ──────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_rect(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(1.0), fill_color=TERRACOTTA)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(0.5), Inches(0.4), height=Inches(0.8))
    add_text(slide, "Sylvan Hills x GBAutomation", Inches(0.5), Inches(1.6),
             Inches(12), Inches(1.2), size=44, bold=True, color=WHITE)
    add_text(slide, "Engagement Roadmap", Inches(0.5), Inches(2.8),
             Inches(10), Inches(0.7), size=24, color=TERRACOTTA)
    add_text(slide, "April 2026", Inches(0.5), Inches(3.8),
             Inches(6), Inches(0.5), size=16, color=RGBColor(0x99, 0x99, 0x99))
    add_text(slide, '"Garments for the man who has nothing to prove."',
             Inches(0.5), Inches(6.65), Inches(6), Inches(0.4),
             size=13, color=WHITE, align=PP_ALIGN.LEFT)


def slide_phase(prs, phase_num, phase_name, status, description, deliverables, timeline):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, f"PHASE {phase_num}")
    add_logo(slide)

    # Phase title
    add_text(slide, f"Phase {phase_num}: {phase_name}",
             Inches(0.4), Inches(0.22), Inches(11.5), Inches(0.6),
             size=26, bold=True, color=DARK)

    # Status badge
    badge_color = TERRACOTTA if status == "ACTIVE" else MID
    add_text(slide, status, Inches(10.5), Inches(0.28), Inches(2), Inches(0.4),
             size=12, bold=True, color=badge_color, align=PP_ALIGN.RIGHT)

    divider(slide, Inches(0.9))

    # Description
    add_text(slide, description, Inches(0.4), Inches(1.1), Inches(12), Inches(0.8),
             size=15, color=MID)

    # Timeline
    add_text(slide, f"Timeline: {timeline}", Inches(0.4), Inches(1.9), Inches(6), Inches(0.4),
             size=13, bold=True, color=TERRA_DARK)

    # Deliverables
    add_text(slide, "Deliverables", Inches(0.4), Inches(2.5), Inches(6), Inches(0.4),
             size=16, bold=True, color=DARK)
    bullets(slide, deliverables, Inches(0.5), Inches(3.0), Inches(11.5), Inches(4.0), size=14)


def slide_requirements(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "REQUIREMENTS")
    add_logo(slide)

    add_text(slide, "What We Need From You",
             Inches(0.4), Inches(0.22), Inches(11.5), Inches(0.6),
             size=26, bold=True, color=DARK)
    divider(slide, Inches(0.9))

    add_text(slide, "To unlock the next phases, we need a few things set up on your end:",
             Inches(0.4), Inches(1.1), Inches(12), Inches(0.5),
             size=15, color=MID)

    items = [
        "GitHub account username — so we can add you as a collaborator on the Sylvan Hills repo",
        (1, "You'll get read access to workspace files, skill definitions, and brief configs"),
        "Instagram Business credentials — for deeper analytics access in Phase 2",
        (1, "Business account API token or Meta Business Suite access"),
        "Competitor list — 3-5 Instagram accounts you watch closely",
        (1, "We'll set up automated tracking for their posting patterns and engagement"),
        "Preferred notification channel — Telegram, WhatsApp, or email for brief delivery",
        "Content calendar access — if you have one (spreadsheet, Notion, etc.)",
        (1, "Helps us understand your posting rhythm for Phase 3"),
    ]
    bullets(slide, items, Inches(0.5), Inches(1.8), Inches(11.5), Inches(5.0), size=14)


def slide_already_set_up(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "INFRASTRUCTURE")
    add_logo(slide)

    add_text(slide, "What's Already Running",
             Inches(0.4), Inches(0.22), Inches(11.5), Inches(0.6),
             size=26, bold=True, color=DARK)
    divider(slide, Inches(0.9))

    items = [
        "Daily Community Brief — cron job running at 2pm PT via Sebastian (AI agent)",
        (1, "Generates Instagram engagement intel and draft comments"),
        (1, "Enhancement incoming: post date + image description per comment"),
        "Google Drive — shared folder with all onboarding docs",
        (1, "drive.google.com/drive/folders/1Czsj7LOeFEwr1DaZvw93yls-jG-2YN1x"),
        "GitHub Repository — gbauto/sylvan-hills (private repo)",
        (1, "Workspace files, skill definitions, agent configuration"),
        "Sebastian — dedicated AI agent on Mac Mini infrastructure",
        (1, "Runs scheduled tasks, manages briefs, powers automation"),
        "Branded onboarding materials — welcome docs, service agreement, prep guide",
    ]
    bullets(slide, items, Inches(0.5), Inches(1.1), Inches(11.5), Inches(5.5), size=14)


def slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_rect(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(1.0), fill_color=TERRACOTTA)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(0.5), Inches(0.3), height=Inches(0.6))

    add_text(slide, "Next Steps", Inches(0.5), Inches(1.2),
             Inches(12), Inches(0.8), size=36, bold=True, color=WHITE)

    items = [
        "Reply to this email with your GitHub username",
        "Share your top 3-5 competitor Instagram accounts",
        "Let us know your preferred notification channel (Telegram / WhatsApp / email)",
        "We'll grant you repo access and start scoping Phase 2",
    ]

    txb = slide.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(11), Inches(4.0))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"{i+1}.  {item}"
        run.font.size = Pt(18)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"
        p.space_after = Pt(12)

    add_text(slide, "greg@gbautomation.xyz  |  gbautomation.xyz",
             Inches(0.5), Inches(6.65), Inches(6), Inches(0.4),
             size=13, color=RGBColor(0x99, 0x99, 0x99))


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title
    slide_title(prs)

    # 2. Phase 1 — Community Briefs
    slide_phase(prs, 1, "Community Briefs", "ACTIVE",
        "Daily Instagram engagement intelligence — identifying high-value posts in your niche "
        "and generating personalized draft comments to build authentic connections with your community.",
        [
            "Daily community brief delivered at 2pm PT",
            "Draft comments tailored to Sylvan Hills brand voice",
            "Post metadata: date + image description per comment (shipping now)",
            "Engagement tracking across target accounts",
        ],
        "Now — Already Running")

    # 3. Phase 2 — Competitor Tracking
    slide_phase(prs, 2, "Competitor Tracking", "NEXT",
        "Automated monitoring of competitor Instagram accounts — posting frequency, engagement rates, "
        "trending aesthetics, hashtag strategy, and audience overlap analysis.",
        [
            "Weekly competitor report across 3-5 tracked accounts",
            "Engagement rate benchmarking (yours vs. competitors)",
            "Trending aesthetics and visual style analysis in menswear",
            "Hashtag performance comparison",
            "Audience overlap and growth rate tracking",
        ],
        "Unlocks after Phase 1 requirements met")

    # 4. Phase 3 — Content Calendar
    slide_phase(prs, 3, "Content Calendar", "PLANNED",
        "AI-assisted content planning and scheduling — optimal posting times, hashtag research, "
        "caption generation, and visual content recommendations based on what's working.",
        [
            "AI-generated content calendar (weekly cadence)",
            "Optimal posting time recommendations",
            "Caption drafts aligned with brand voice",
            "Hashtag sets optimized for reach and engagement",
            "Visual content recommendations based on competitor analysis",
        ],
        "Scoping begins after Phase 2 baseline")

    # 5. Phase 4 — Growth Analytics
    slide_phase(prs, 4, "Growth Analytics", "PLANNED",
        "Comprehensive growth dashboards and audience insights — follower acquisition, "
        "engagement trends, content performance scoring, and strategic growth recommendations.",
        [
            "Growth dashboard with key metrics",
            "Follower acquisition and churn tracking",
            "Content performance scoring (what works, what doesn't)",
            "Audience demographic insights",
            "Monthly strategic growth recommendations",
        ],
        "Builds on Phase 2-3 data foundation")

    # 6. Requirements
    slide_requirements(prs)

    # 7. What's Already Set Up
    slide_already_set_up(prs)

    # 8. Next Steps
    slide_next_steps(prs)

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
