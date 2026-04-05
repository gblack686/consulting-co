"""
Generate a branded GBAutomation client onboarding PowerPoint deck.

Usage:
    python -m scripts.generate_client_deck \
        --name "Erica Cruz" \
        --email "artbycruzcreations@gmail.com" \
        --session-date "Thursday, March 5, 2026" \
        --session-time "4:00 PM" \
        --timezone "PST" \
        --meet-link "https://meet.google.com/vjk-dvpq-nfe" \
        --output "erica-cruz-onboarding.pptx"
"""

import argparse
import sys
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy

# ─── Brand Colors ────────────────────────────────────────────────────────────
CREAM      = RGBColor(0xF3, 0xF1, 0xE7)
TERRACOTTA = RGBColor(0xD9, 0x77, 0x57)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)
MID        = RGBColor(0x55, 0x55, 0x55)
LIGHT_TAN  = RGBColor(0xE8, 0xE5, 0xD8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TERRA_DARK = RGBColor(0xB5, 0x5E, 0x3A)

SKILL_ROOT  = Path(__file__).parent.parent
LOGO_PATH   = SKILL_ROOT / "assets" / "gb-logo.png"
OUTPUT_DIR  = SKILL_ROOT / "output"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── Helpers ─────────────────────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb

def add_logo(slide, right_margin=Inches(0.3), top=Inches(0.15), height=Inches(0.55)):
    if LOGO_PATH.exists():
        aspect = 1.0
        try:
            from PIL import Image
            with Image.open(LOGO_PATH) as img:
                w, h = img.size
                aspect = w / h
        except Exception:
            aspect = 3.5
        width = height * aspect
        left = SLIDE_W - width - right_margin
        slide.shapes.add_picture(str(LOGO_PATH), left, top, width, height)

def section_bar(slide, label: str):
    """Terracotta left bar with section label."""
    bar = add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill_color=TERRACOTTA)
    add_text(slide, label, Inches(0.28), Inches(0.18), Inches(4), Inches(0.45),
             size=11, bold=True, color=TERRACOTTA)

def divider(slide, top, color=LIGHT_TAN, width_in=12.5):
    add_rect(slide, Inches(0.4), top, Inches(width_in), Pt(1.5), fill_color=color)

def bullet_frame(slide, items, l, t, w, h, size=14, level_indent=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        indent = "   " * level if level_indent else ""
        run = p.add_run()
        run.text = ("• " if level == 0 else "◦ ") + indent + text
        run.font.size = Pt(size - level * 1.5)
        run.font.color.rgb = DARK if level == 0 else MID
        run.font.name = "Calibri"
    return txb

def table_slide(prs, section, title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, CREAM)
    section_bar(slide, section)
    add_logo(slide)
    # Title
    add_text(slide, title, Inches(0.4), Inches(0.22), Inches(11.5), Inches(0.6),
             size=22, bold=True, color=DARK)
    divider(slide, Inches(0.9))

    ncols = len(headers)
    nrows = len(rows) + 1
    tbl_l = Inches(0.4)
    tbl_t = Inches(1.0)
    tbl_w = Inches(12.5)
    tbl_h = Inches(0.4 * nrows + 0.2)

    tbl = slide.shapes.add_table(nrows, ncols, tbl_l, tbl_t, tbl_w, tbl_h).table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TERRACOTTA
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_TAN
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(11)
            run.font.color.rgb = DARK
            run.font.name = "Calibri"

    add_logo(slide)
    return slide


# ─── Slide Builders ───────────────────────────────────────────────────────────

def slide_title(prs, client_name, session_date, session_time, timezone):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)

    # Big terracotta accent bar bottom
    add_rect(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(1.0), fill_color=TERRACOTTA)

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(0.5), Inches(0.4), height=Inches(0.8))

    add_text(slide, f"Welcome, {client_name.split()[0]}.", Inches(0.5), Inches(1.6),
             Inches(12), Inches(1.2), size=44, bold=True, color=WHITE)

    add_text(slide, "Your AI Agent Setup — Everything You Need to Know",
             Inches(0.5), Inches(2.8), Inches(10), Inches(0.7),
             size=20, color=RGBColor(0xCC, 0xC8, 0xBC))

    add_text(slide, f"{session_date}  ·  {session_time} {timezone}",
             Inches(0.5), Inches(3.6), Inches(9), Inches(0.55),
             size=16, color=TERRACOTTA, bold=True)

    add_text(slide, "GBAutomation  ·  gbautomation.xyz",
             Inches(0.5), Inches(6.6), Inches(8), Inches(0.5),
             size=13, color=WHITE)
    return slide


def slide_what_to_expect(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "WELCOME")
    add_logo(slide)
    add_text(slide, "What to Expect", Inches(0.4), Inches(0.22), Inches(11), Inches(0.6),
             size=26, bold=True, color=DARK)
    divider(slide, Inches(0.95))

    phases = [
        ("BEFORE THE CALL", TERRACOTTA, [
            "Complete the Pre-Session Prep Guide (~15 min)",
            "List your daily tools and apps",
            "Think through your top 3–5 work domains",
            "Identify your most time-consuming recurring task",
        ]),
        ("DURING THE CALL  (60 MIN)", DARK, [
            "Walk through your goals, tools, workflows, and preferences",
            "We'll define your agent's personality and boundaries",
            "Map out your top automation opportunities",
        ]),
        ("AFTER THE CALL", TERRACOTTA, [
            "We build your full agent system within 5 business days",
            "Deployment + walkthrough call included",
            "You own everything — all files, skills, and configs",
        ]),
    ]

    col_l = [Inches(0.35), Inches(4.55), Inches(8.75)]
    col_w = Inches(3.9)
    for i, (label, col, bullets) in enumerate(phases):
        l = col_l[i]
        # Header card
        card = add_rect(slide, l, Inches(1.1), col_w, Inches(0.5), fill_color=col)
        add_text(slide, label, l + Inches(0.12), Inches(1.15), col_w - Inches(0.2), Inches(0.4),
                 size=11, bold=True, color=WHITE)
        # Bullet area
        add_rect(slide, l, Inches(1.6), col_w, Inches(5.3), fill_color=WHITE)
        bullet_frame(slide, bullets, l + Inches(0.15), Inches(1.75),
                     col_w - Inches(0.25), Inches(5.0), size=13)
    return slide


def slide_deliverables(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "WELCOME")
    add_logo(slide)
    add_text(slide, "What You'll Walk Away With", Inches(0.4), Inches(0.22),
             Inches(11), Inches(0.6), size=26, bold=True, color=DARK)
    divider(slide, Inches(0.95))

    items = [
        ("Fully Configured AI Agent", "Running on your own infrastructure — your device, your cloud"),
        ("Custom Skills", "Step-by-step automations for your top recurring workflows"),
        ("Domain Expert Systems", "Self-improving knowledge bases for each area of your work"),
        ("Scheduled Automation", "Cron jobs for morning briefs, reports, and recurring tasks"),
        ("Complete Documentation", "Every component explained — so you can grow it yourself"),
        ("Walkthrough Call", "Live demo showing exactly what your agent can do"),
    ]

    for i, (title, desc) in enumerate(items):
        row = i // 2
        col = i % 2
        l = Inches(0.4) + col * Inches(6.4)
        t = Inches(1.1) + row * Inches(1.8)
        add_rect(slide, l, t, Inches(6.1), Inches(1.6), fill_color=WHITE)
        # Terracotta left accent
        add_rect(slide, l, t, Inches(0.06), Inches(1.6), fill_color=TERRACOTTA)
        add_text(slide, title, l + Inches(0.18), t + Inches(0.15),
                 Inches(5.7), Inches(0.45), size=14, bold=True, color=DARK)
        add_text(slide, desc, l + Inches(0.18), t + Inches(0.6),
                 Inches(5.7), Inches(0.9), size=12, color=MID)
    return slide


def slide_checklist(prs, meet_link):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "BEFORE THE CALL")
    add_logo(slide)
    add_text(slide, "Have These Ready", Inches(0.4), Inches(0.22),
             Inches(11), Inches(0.6), size=26, bold=True, color=DARK)
    divider(slide, Inches(0.95))

    checklist = [
        "List of apps & tools you use daily (Gmail, Notion, Slack, etc.)",
        "Any existing API keys (Anthropic, OpenAI, Brave, Stripe, etc.)",
        "Access to the device where your agent will run (Mac Mini, VPS, laptop)",
        "Your phone for WhatsApp or Telegram channel setup",
        "Rough idea of your 3–5 work 'departments'",
        "The most annoying recurring task you'd love to never do again",
    ]

    for i, item in enumerate(checklist):
        t = Inches(1.1) + i * Inches(0.82)
        # Checkbox
        add_rect(slide, Inches(0.4), t + Inches(0.05), Inches(0.38), Inches(0.38),
                 line_color=TERRACOTTA, line_width=Pt(2))
        add_text(slide, item, Inches(0.9), t, Inches(11.5), Inches(0.75),
                 size=14, color=DARK)

    divider(slide, Inches(6.2))
    add_text(slide, f"Session link: {meet_link}",
             Inches(0.4), Inches(6.3), Inches(12), Inches(0.5),
             size=12, color=MID, italic=True)
    return slide


def slide_prep_voice(prs, meet_link):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "PRE-SESSION PREP")
    add_logo(slide)
    add_text(slide, "Prefer to Talk? Use Voice Intake Instead",
             Inches(0.4), Inches(0.22), Inches(11), Inches(0.6),
             size=26, bold=True, color=DARK)
    divider(slide, Inches(0.95))

    add_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(1.7), fill_color=TERRACOTTA)
    add_text(slide, "Join this Meet link and speak your answers — it auto-transcribes, nothing gets lost.",
             Inches(0.65), Inches(1.15), Inches(12.0), Inches(0.5), size=14, color=WHITE)
    add_text(slide, meet_link, Inches(0.65), Inches(1.65), Inches(12.0), Inches(0.5),
             size=16, bold=True, color=WHITE)

    questions = [
        ("1", "Your Tools", "Walk through every app you use at least weekly. What do you use it for?"),
        ("2", "Your Departments", "If you hired 3–5 people to replace you, what would each person own?"),
        ("3", "Your Most Annoying Task", "One recurring task that takes 30+ min and follows the same steps every time."),
        ("4", "Perfect Morning Briefing", "What would a useful daily AI message tell you? Data, format, tone?"),
        ("5", "Agent Personality", "What should your agent call you? Professional, casual, playful?"),
    ]

    for i, (num, title, desc) in enumerate(questions):
        col = i % 3
        row = i // 3
        l = Inches(0.35) + col * Inches(4.3)
        t = Inches(3.0) + row * Inches(2.1)
        add_rect(slide, l, t, Inches(4.05), Inches(1.95), fill_color=WHITE)
        add_rect(slide, l, t, Inches(0.38), Inches(0.38), fill_color=TERRACOTTA)
        add_text(slide, num, l + Inches(0.07), t + Inches(0.03), Inches(0.3), Inches(0.35),
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, l + Inches(0.45), t + Inches(0.05), Inches(3.5), Inches(0.4),
                 size=13, bold=True, color=DARK)
        add_text(slide, desc, l + Inches(0.12), t + Inches(0.55), Inches(3.8), Inches(1.3),
                 size=11, color=MID)
    return slide


def slide_tool_inventory(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "PRE-SESSION PREP")
    add_logo(slide)
    add_text(slide, "Exercise 1 — Your Tool Inventory",
             Inches(0.4), Inches(0.22), Inches(11), Inches(0.6),
             size=26, bold=True, color=DARK)
    add_text(slide, "List every app or tool you use at least weekly. Don't know if it has an API? Just list it.",
             Inches(0.4), Inches(0.85), Inches(12), Inches(0.4), size=13, color=MID)
    divider(slide, Inches(1.3))

    headers = ["App / Tool", "What You Use It For", "Has an API?", "Have API Key?"]
    rows = [["", "", "Yes / No / ?", "Yes / No"] for _ in range(7)]
    rows[0] = ["Example: Notion", "Project management, notes", "Yes", "No"]

    tbl_l, tbl_t = Inches(0.4), Inches(1.4)
    tbl_w, tbl_h = Inches(12.5), Inches(5.7)
    ncols, nrows = 4, 8
    tbl = slide.shapes.add_table(nrows, ncols, tbl_l, tbl_t, tbl_w, tbl_h).table

    col_widths_in = [3.2, 5.0, 2.2, 2.1]
    for i, cw in enumerate(col_widths_in):
        tbl.columns[i].width = Inches(cw)

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TERRACOTTA
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_TAN if ri == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(11)
            run.font.color.rgb = MID if ri == 0 else DARK
            run.font.italic = (ri == 0)
            run.font.name = "Calibri"
    return slide


def slide_departments(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "PRE-SESSION PREP")
    add_logo(slide)
    add_text(slide, "Exercise 2 — Your Departments",
             Inches(0.4), Inches(0.22), Inches(11), Inches(0.6),
             size=26, bold=True, color=DARK)
    add_text(slide, "If you hired 3–5 people to replace you, what would each person own? Give each role a name.",
             Inches(0.4), Inches(0.85), Inches(12), Inches(0.4), size=13, color=MID)
    divider(slide, Inches(1.3))

    examples = ["Content Production", "Business Development", "Client Ops", "Personal Admin", ""]
    for i in range(5):
        t = Inches(1.5) + i * Inches(1.05)
        add_rect(slide, Inches(0.4), t, Inches(0.45), Inches(0.45), fill_color=TERRACOTTA)
        add_text(slide, str(i + 1), Inches(0.4), t + Inches(0.04), Inches(0.45), Inches(0.4),
                 size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(0.95), t, Inches(7.5), Inches(0.45),
                 fill_color=WHITE, line_color=LIGHT_TAN, line_width=Pt(1))
        if examples[i]:
            add_text(slide, 'e.g. \u201c' + examples[i] + '\u201d', Inches(1.1), t + Inches(0.05),
                     Inches(7.2), Inches(0.38), size=12, color=LIGHT_TAN, italic=True)
        add_text(slide, "→ Owns:", Inches(8.6), t + Inches(0.05),
                 Inches(0.8), Inches(0.38), size=11, color=MID)
        add_rect(slide, Inches(9.45), t, Inches(3.4), Inches(0.45),
                 fill_color=WHITE, line_color=LIGHT_TAN, line_width=Pt(1))
    return slide


def slide_annoying_task(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "PRE-SESSION PREP")
    add_logo(slide)
    add_text(slide, "Exercise 3 — Your Most Annoying Task",
             Inches(0.4), Inches(0.22), Inches(11), Inches(0.6),
             size=26, bold=True, color=DARK)
    add_text(slide, "The one recurring task you'd love to automate away forever.",
             Inches(0.4), Inches(0.85), Inches(12), Inches(0.4), size=13, color=MID)
    divider(slide, Inches(1.3))

    fields = [
        ("The task", Inches(5.5), Inches(0.55)),
        ("How often", Inches(3.0), Inches(0.55)),
    ]
    t_start = Inches(1.45)
    for label, fw, fh in fields:
        add_text(slide, label, Inches(0.4), t_start, Inches(2.5), Inches(0.4),
                 size=13, bold=True, color=DARK)
        add_rect(slide, Inches(2.9), t_start - Inches(0.03), fw, fh,
                 fill_color=WHITE, line_color=LIGHT_TAN, line_width=Pt(1))
        t_start += Inches(0.75)

    add_text(slide, "Rough steps (bullet by bullet):",
             Inches(0.4), Inches(3.05), Inches(12), Inches(0.4),
             size=13, bold=True, color=DARK)

    for i in range(5):
        t = Inches(3.55) + i * Inches(0.72)
        add_rect(slide, Inches(0.4), t + Inches(0.12), Inches(0.2), Inches(0.2),
                 fill_color=TERRACOTTA)
        add_rect(slide, Inches(0.75), t, Inches(11.8), Inches(0.55),
                 fill_color=WHITE, line_color=LIGHT_TAN, line_width=Pt(1))

    add_rect(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3), fill_color=TERRACOTTA)
    add_text(slide, "⚡  This is usually the first skill we build for you.",
             Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
             size=11, color=WHITE, bold=True)
    return slide


def slide_morning_briefing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "PRE-SESSION PREP")
    add_logo(slide)
    add_text(slide, "Exercise 4 — Your Perfect Morning Briefing",
             Inches(0.4), Inches(0.22), Inches(11.5), Inches(0.6),
             size=26, bold=True, color=DARK)
    add_text(slide, "What would make you say \"this message is incredibly useful\" every morning?",
             Inches(0.4), Inches(0.85), Inches(12), Inches(0.4), size=13, color=MID)
    divider(slide, Inches(1.3))

    options = [
        "Weather for today",
        "Calendar summary",
        "Priority emails / messages",
        "News in my industry / niche",
        "Social media analytics",
        "Task list / to-dos",
        "Project status updates",
        "Financial summary (revenue, expenses)",
        "Content performance (views, engagement)",
        "Custom metric or data feed",
        "Other: _______________",
    ]

    ncols = 2
    col_w = Inches(6.1)
    for i, opt in enumerate(options):
        col = i % ncols
        row = i // ncols
        l = Inches(0.4) + col * Inches(6.5)
        t = Inches(1.45) + row * Inches(0.6)
        add_rect(slide, l, t + Inches(0.08), Inches(0.3), Inches(0.3),
                 line_color=TERRACOTTA, line_width=Pt(2))
        add_text(slide, opt, l + Inches(0.45), t, col_w, Inches(0.55),
                 size=13, color=DARK)
    return slide


def slide_agent_personality(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "PRE-SESSION PREP")
    add_logo(slide)
    add_text(slide, "Bonus — Your Agent's Personality",
             Inches(0.4), Inches(0.22), Inches(11), Inches(0.6),
             size=26, bold=True, color=DARK)
    add_text(slide, "No pressure — we'll brainstorm on the call. Just start thinking.",
             Inches(0.4), Inches(0.85), Inches(12), Inches(0.4), size=13, color=MID)
    divider(slide, Inches(1.3))

    questions = [
        ("What should your agent call you?",
         "First name? Nickname? Title? (e.g. Greg, G, Boss)",
         Inches(1.45)),
        ("What vibe do you want?",
         "Professional assistant  ·  Casual friend  ·  Sharp executive  ·  Playful",
         Inches(2.95)),
        ("Got a name in mind for your agent?",
         "Human name, creative name, acronym — or we'll brainstorm together",
         Inches(4.45)),
    ]

    for label, hint, t in questions:
        add_text(slide, label, Inches(0.4), t, Inches(5.5), Inches(0.45),
                 size=14, bold=True, color=DARK)
        add_text(slide, hint, Inches(0.4), t + Inches(0.5), Inches(5.5), Inches(0.45),
                 size=11, color=MID, italic=True)
        add_rect(slide, Inches(6.2), t, Inches(6.7), Inches(0.8),
                 fill_color=WHITE, line_color=LIGHT_TAN, line_width=Pt(1))

    # Vibe options
    vibes = ["Professional", "Casual", "Sharp", "Playful", "Custom"]
    for i, v in enumerate(vibes):
        l = Inches(6.3) + i * Inches(1.3)
        add_rect(slide, l, Inches(3.1), Inches(1.15), Inches(0.42),
                 fill_color=LIGHT_TAN if v != "Custom" else WHITE,
                 line_color=TERRACOTTA, line_width=Pt(1))
        add_text(slide, v, l + Inches(0.05), Inches(3.13), Inches(1.05), Inches(0.38),
                 size=11, color=DARK, align=PP_ALIGN.CENTER)
    return slide


def slide_agenda_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "SESSION AGENDA")
    add_logo(slide)
    add_text(slide, "60-Minute Session — What We'll Cover",
             Inches(0.4), Inches(0.22), Inches(11), Inches(0.6),
             size=26, bold=True, color=DARK)
    divider(slide, Inches(0.95))

    phases = [
        ("Phase 1", "Identity & Soul", "15 min", "Define who your agent is"),
        ("Phase 2", "Tools & Infrastructure", "15 min", "Map your tech stack"),
        ("Phase 3", "Domain Discovery", "30 min", "Map your top workflows"),
        ("Phase 4", "Autonomy & Safety", "10 min", "Set permissions and limits"),
        ("Phase 5", "Synthesis", "10 min", "Review and confirm everything"),
    ]

    for i, (phase, title, duration, desc) in enumerate(phases):
        t = Inches(1.1) + i * Inches(1.18)
        # Phase number pill
        add_rect(slide, Inches(0.4), t, Inches(1.1), Inches(0.5),
                 fill_color=TERRACOTTA)
        add_text(slide, phase, Inches(0.4), t + Inches(0.05), Inches(1.1), Inches(0.42),
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, Inches(1.6), t, Inches(6.0), Inches(0.5),
                 size=15, bold=True, color=DARK)
        add_text(slide, desc, Inches(1.6), t + Inches(0.45), Inches(6.0), Inches(0.4),
                 size=12, color=MID)
        add_rect(slide, Inches(9.8), t, Inches(1.8), Inches(0.5), fill_color=LIGHT_TAN)
        add_text(slide, duration, Inches(9.8), t + Inches(0.05), Inches(1.8), Inches(0.42),
                 size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Progress bar
    add_rect(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.22), fill_color=LIGHT_TAN)
    add_rect(slide, Inches(0.4), Inches(7.05), Inches(12.5 * (60/90)), Inches(0.22), fill_color=TERRACOTTA)
    return slide


def slide_after_session(prs):
    headers = ["Timeline", "What Happens"]
    rows = [
        ["Day 1", "Process your transcript — build workspace files (SOUL.md, USER.md, etc.)"],
        ["Days 2–3", "Build domain expert systems — one per department"],
        ["Days 3–4", "Quality validation and testing"],
        ["Days 4–5", "Deployment to your OpenClaw instance"],
        ["Day 5–7", "Walkthrough call — live demo of your full agent"],
    ]
    slide = table_slide(prs, "AFTER THE SESSION", "Delivery Timeline",
                        headers, rows, col_widths=[2.5, 10.0])
    add_rect(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.7), fill_color=TERRACOTTA)
    add_text(slide, "You'll receive: all workspace files  ·  expert systems  ·  deployed skills  ·  quality report  ·  walkthrough call",
             Inches(0.6), Inches(6.58), Inches(12.0), Inches(0.5),
             size=12, bold=True, color=WHITE)
    return slide


def slide_key_terms_agent(prs):
    headers = ["Term", "What It Means"]
    rows = [
        ["OpenClaw", "The platform that runs your AI agent 24/7 — the OS for your personal AI"],
        ["Agent", "Your AI assistant. Has a name, personality, memory, and skills"],
        ["Gateway", "Keeps your agent online and connected to messaging channels"],
        ["Workspace", "The folder where all your agent's configuration lives"],
        ["SOUL.md", "Your agent's values and communication style — like a constitution"],
        ["MEMORY.md", "Long-term facts your agent should never forget"],
        ["Skill", "A specific task your agent knows how to do — like a recipe"],
        ["Cron job", "A scheduled task that runs automatically (e.g. morning briefing at 7am)"],
    ]
    return table_slide(prs, "KEY TERMS", "Your Agent — Quick Reference",
                       headers, rows, col_widths=[2.8, 9.7])


def slide_key_terms_autonomy(prs):
    headers = ["Term", "What It Means"]
    rows = [
        ["Domain", "A 'department' of your work (e.g. Content, Business, Personal)"],
        ["Workflow", "A repeatable process within a domain"],
        ["Trigger", "What starts a workflow: time, event, heartbeat, or on-demand request"],
        ["Heartbeat", "Periodic check-in — agent wakes up every 30 min to run a checklist"],
        ["Approval gate", "A point where the agent stops and asks you before continuing"],
        ["Blast radius", "What happens if something goes wrong — determines autonomy level"],
        ["Autonomy level", "How independent your agent is: ask-everything ↔ just-get-it-done"],
        ["Expert system", "Self-improving knowledge package for one domain — gets smarter over time"],
    ]
    return table_slide(prs, "KEY TERMS", "Skills, Automation & Autonomy",
                       headers, rows, col_widths=[2.8, 9.7])


def slide_models(prs):
    headers = ["Model", "Speed", "Cost", "Best For"]
    rows = [
        ["Haiku", "Fast", "$", "Simple tasks — file moves, data extraction, formatting"],
        ["Sonnet", "Balanced", "$$", "Most tasks — writing, research, scheduling, analysis"],
        ["Opus", "Thorough", "$$$", "Complex judgment — strategy, quality review, edge cases"],
    ]
    return table_slide(prs, "KEY TERMS", "AI Models — Quick Reference",
                       headers, rows, col_widths=[2.5, 2.0, 1.5, 6.5])


def slide_pricing(prs, foundation_price="$1,500", standard_price="$2,500", premium_price="$4,000",
                  stripe_link="https://gbautomation.xyz/pay"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)
    section_bar(slide, "SERVICE AGREEMENT")
    add_logo(slide)
    add_text(slide, "Service Tiers & Pricing",
             Inches(0.4), Inches(0.22), Inches(11), Inches(0.6),
             size=26, bold=True, color=DARK)
    divider(slide, Inches(0.95))

    tiers = [
        ("Foundation", "1–2 Domains", "Up to 4 skills", foundation_price, False),
        ("Standard", "3–4 Domains", "Up to 8 skills", standard_price, True),
        ("Premium", "5+ Domains", "Unlimited skills", premium_price, False),
    ]

    for i, (name, domains, skills, price, featured) in enumerate(tiers):
        l = Inches(0.4) + i * Inches(4.3)
        bg = TERRACOTTA if featured else WHITE
        fg = WHITE if featured else DARK
        sfg = WHITE if featured else MID

        add_rect(slide, l, Inches(1.1), Inches(4.0), Inches(5.4), fill_color=bg,
                 line_color=TERRACOTTA if not featured else None, line_width=Pt(1.5))
        if featured:
            add_text(slide, "MOST POPULAR", l, Inches(0.85), Inches(4.0), Inches(0.3),
                     size=9, bold=True, color=TERRACOTTA, align=PP_ALIGN.CENTER)

        add_text(slide, name, l + Inches(0.2), Inches(1.25), Inches(3.6), Inches(0.55),
                 size=20, bold=True, color=fg)
        add_text(slide, price, l + Inches(0.2), Inches(1.9), Inches(3.6), Inches(0.65),
                 size=28, bold=True, color=fg)
        add_text(slide, domains, l + Inches(0.2), Inches(2.7), Inches(3.6), Inches(0.45),
                 size=13, color=sfg)
        add_text(slide, skills, l + Inches(0.2), Inches(3.2), Inches(3.6), Inches(0.45),
                 size=13, color=sfg)

        incl = ["Discovery session (60 min)", "Workspace files + gateway config",
                "Domain expert systems", "Deployment + walkthrough call", "Quality report"]
        for j, inc in enumerate(incl):
            add_text(slide, "✓  " + inc, l + Inches(0.2), Inches(3.85) + j * Inches(0.42),
                     Inches(3.6), Inches(0.38), size=11, color=sfg)

    add_rect(slide, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.42), fill_color=LIGHT_TAN)
    add_text(slide, f"Payment via Stripe: {stripe_link}  ·  Visa, Mastercard, Amex accepted",
             Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.35),
             size=11, color=MID, italic=True)
    return slide


def slide_next_steps(prs, client_name, session_date, session_time, meet_link, stripe_link="https://gbautomation.xyz/pay"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_logo(slide)

    add_text(slide, "You're All Set.", Inches(0.5), Inches(0.6),
             Inches(12), Inches(0.9), size=36, bold=True, color=WHITE)

    steps = [
        ("1", "Complete the Pre-Session Prep", "Exercises 1–4 — takes about 15 minutes"),
        ("2", "Sign the Service Agreement", "In your Google Drive folder — link sent to your email"),
        ("3", "Complete Payment", f"Stripe: {stripe_link}"),
        ("4", "Join on Session Day", f"{session_date}  ·  {session_time} PST  ·  {meet_link}"),
    ]

    for i, (num, title, sub) in enumerate(steps):
        t = Inches(1.65) + i * Inches(1.28)
        add_rect(slide, Inches(0.5), t, Inches(0.55), Inches(0.55), fill_color=TERRACOTTA)
        add_text(slide, num, Inches(0.5), t + Inches(0.06), Inches(0.55), Inches(0.45),
                 size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, Inches(1.2), t, Inches(11), Inches(0.5),
                 size=16, bold=True, color=WHITE)
        add_text(slide, sub, Inches(1.2), t + Inches(0.5), Inches(11), Inches(0.5),
                 size=12, color=RGBColor(0xCC, 0xC8, 0xBC))

    add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), fill_color=TERRACOTTA)
    add_text(slide, "Questions? Reply to your welcome email or message greg@gbautomation.xyz",
             Inches(0.5), Inches(6.88), Inches(12), Inches(0.45),
             size=12, color=WHITE)
    return slide


# ─── Main ─────────────────────────────────────────────────────────────────────

def generate(args):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_title(prs, args.name, args.session_date, args.session_time, args.timezone)
    slide_what_to_expect(prs)
    slide_deliverables(prs)
    slide_checklist(prs, args.meet_link)
    slide_prep_voice(prs, args.meet_link)
    slide_tool_inventory(prs)
    slide_departments(prs)
    slide_annoying_task(prs)
    slide_morning_briefing(prs)
    slide_agent_personality(prs)
    slide_agenda_overview(prs)
    slide_after_session(prs)
    slide_key_terms_agent(prs)
    slide_key_terms_autonomy(prs)
    slide_models(prs)
    slide_pricing(prs)
    slide_next_steps(prs, args.name, args.session_date, args.session_time, args.meet_link)

    OUTPUT_DIR.mkdir(exist_ok=True)
    out = OUTPUT_DIR / args.output
    prs.save(str(out))
    print(f"\n✓ Deck saved: {out}")
    print(f"  Slides: {len(prs.slides)}")
    return out


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--name",         default="Erica Cruz")
    parser.add_argument("--email",        default="artbycruzcreations@gmail.com")
    parser.add_argument("--session-date", default="Thursday, March 5, 2026")
    parser.add_argument("--session-time", default="4:00 PM")
    parser.add_argument("--timezone",     default="PST")
    parser.add_argument("--meet-link",    default="https://meet.google.com/vjk-dvpq-nfe")
    parser.add_argument("--output",       default="erica-cruz-onboarding.pptx")
    args = parser.parse_args()
    generate(args)
