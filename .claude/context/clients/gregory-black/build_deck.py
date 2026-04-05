"""Build Gregory Black Client Research PowerPoint Deck"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CREAM = RGBColor(0xF3, 0xF1, 0xE7)
TERRACOTTA = RGBColor(0xD9, 0x77, 0x57)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=CREAM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_bullet(tf, text, size=16, color=DARK, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(4)
    return p

def accent_bar(slide, top=1.3):
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(top), Inches(2), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TERRACOTTA
    shape.line.fill.background()

# === SLIDE 1: Title ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 1.5, 11, 1.2, "Gregory Black", size=48, bold=True, color=DARK)
accent_bar(s, top=2.8)
add_text(s, 0.8, 3.1, 11, 1.5,
    "Client Research Dossier\nData & Automation Expert | Agentic Engineering | Claude Code SDK",
    size=22, color=GRAY)
add_text(s, 0.8, 5.5, 6, 0.5, "GBAutomation Consulting  |  March 8, 2026", size=14, color=GRAY)

# === SLIDE 2: Profile Overview ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Profile Overview", size=36, bold=True, color=DARK)
accent_bar(s)

rows = [
    ("Name", "Gregory Black"),
    ("Headline", "Data & Automation Expert | Agentic Engineering | Production-Ready AI Pipelines | Claude Code SDK Implementations"),
    ("Location", "New York, New York (bio says Venice, LA)"),
    ("Current Role", "Lead AI Engineer at Accenture Federal Services (Jan 2026)"),
    ("Side Business", "GB Automation — Claude Code-driven agentic systems (Jun 2024)"),
    ("Education", "Penn State (BS Real Estate Economics) + UC Irvine (Data Analytics)"),
    ("Connections", "500+ LinkedIn | 1,428 IG | ~3M Threads"),
    ("LinkedIn", "linkedin.com/in/gregory-black-gbautomation/"),
    ("Instagram", "@gregorablack | @gb.noir.lugar (secondary)"),
]
tf = add_text(s, 0.8, 1.7, 11.5, 5, "", size=16)
tf.paragraphs[0].text = ""
for label, val in rows:
    add_bullet(tf, f"{label}:  {val}", size=16, bold=False)

# === SLIDE 3: Career Timeline ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Career Timeline", size=36, bold=True, color=DARK)
accent_bar(s)

roles = [
    ("2026 – Present", "Lead AI Engineer", "Accenture Federal Services", "Full-time, 3 months"),
    ("2024 – Present", "Founder / AI Consultant", "GB Automation (Self-employed)", "Claude Code agentic systems, 1yr 10 mos"),
    ("~2021 – 2024", "Sr Data and AI Engineer", "(Company visible in profile)", "ETL, PowerShell, Power Automate, ArcGIS, PBI"),
    ("2019 – 2021", "Data Science Associate", "Axtria - Ingenious Insights", "Pharma consulting, Dataiku DSS, 1B+ records in Redshift"),
    ("2019", "Excel VBA Consulting", "Econsult (Santiago, Chile)", "6 months international consulting"),
    ("2017 – 2018", "Associate", "Province West (LA)", "Land acquisitions, real estate development"),
    ("2015 – 2016", "Analyst", "The Abbey Company (LA)", "CRE operations, 150K sf property mgmt"),
    ("2014", "Capital Markets Intern", "Cassidy Turley (NYC)", "DCF, lease analyses, waterfall structures"),
]
tf = add_text(s, 0.8, 1.7, 11.5, 5.5, "", size=14)
tf.paragraphs[0].text = ""
for dates, title, co, desc in roles:
    add_bullet(tf, f"{dates}  —  {title}  @  {co}", size=15, bold=True)
    add_bullet(tf, f"     {desc}", size=13, color=GRAY)

# === SLIDE 4: Company Deep Dive ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Company Deep Dive", size=36, bold=True, color=DARK)
accent_bar(s)

tf = add_text(s, 0.8, 1.7, 5.5, 5, "", size=16)
tf.paragraphs[0].text = ""
add_bullet(tf, "Accenture Federal Services (AFS)", size=20, bold=True, color=TERRACOTTA)
add_bullet(tf, "Federal IT consulting arm of Accenture", size=15)
add_bullet(tf, "Serves DoD, Intel Community, civilian agencies", size=15)
add_bullet(tf, "~$5.5B revenue, 15,500 employees", size=15)
add_bullet(tf, "#11 Top 100 US Government Contractors", size=15)
add_bullet(tf, "HQ: Arlington, VA | afs.com", size=15)
add_bullet(tf, "Focus: AI/GenAI, cloud, cybersecurity, Salesforce/ServiceNow", size=15)
add_bullet(tf, "CEO: Ron Ash (appointed Sep 2024)", size=15)

tf2 = add_text(s, 6.8, 1.7, 5.5, 5, "", size=16)
tf2.paragraphs[0].text = ""
add_bullet(tf2, "GB Automation (Side Business)", size=20, bold=True, color=TERRACOTTA)
add_bullet(tf2, "gbautomation.xyz (under construction)", size=15)
add_bullet(tf2, "Self-employed since Jun 2024", size=15)
add_bullet(tf2, "25+ n8n/Zapier automations, 5+ Supabase backends", size=15)
add_bullet(tf2, "Clients: UK FinTech, Private Jet Co, Music Mgmt", size=15)
add_bullet(tf2, "Tech: Claude API, LangChain, RAG, LangGraph, OpenRouter", size=15)
add_bullet(tf2, "Booking link in Instagram bio", size=15)

# === SLIDE 5: GB Automation Projects ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "GB Automation — Notable Projects", size=36, bold=True, color=DARK)
accent_bar(s)

projects = [
    ("AI Sales Operations System (FinTech)", "End-to-end sales ops: CRM, outreach automation, performance analytics for UK financial advisory"),
    ("Private Jet Lead Generation", "Automated lead gen + opportunity mgmt with PDF parsing and Salesforce integration"),
    ("SMB Acquisition AI Consultant", "LLM-powered consultation using fine-tuned models from real consultation calls"),
    ("AI SDR Campaign Manager", "Automated lead qualification + SMS nurture with 40% conversion rate"),
    ("Legal CRM Optimization", "Account hierarchy and data unification in Salesforce"),
    ("Automated Job Search System", "AI-powered job matching and outreach via Gmail and LinkedIn"),
    ("Spotify Artist Lead Pipeline", "Automated artist discovery for Joywave Management (Michael Miller testimonial)"),
]
tf = add_text(s, 0.8, 1.7, 11.5, 5.5, "", size=14)
tf.paragraphs[0].text = ""
for name, desc in projects:
    add_bullet(tf, name, size=16, bold=True, color=TERRACOTTA)
    add_bullet(tf, f"  {desc}", size=14, color=GRAY)

# === SLIDE 6: Tech Stack ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Technical Stack", size=36, bold=True, color=DARK)
accent_bar(s)

stacks = [
    ("AI/ML", "LangChain, OpenAI API, Claude API, RAG, Weaviate, Dataiku DSS"),
    ("Automation", "n8n, Zapier, Make.com, Power Automate"),
    ("Cloud", "AWS (Redshift, Lambda), GCP, Supabase"),
    ("CRM", "Salesforce, monday.com"),
    ("Data", "Python, SQL, PostgreSQL, ETL pipelines"),
    ("Web", "FastAPI, Docker, Git"),
    ("BI/Viz", "Tableau, Power BI, ArcGIS"),
    ("Legacy", "Excel VBA, PowerShell, Argus (CRE)"),
]
tf = add_text(s, 0.8, 1.7, 11.5, 5, "", size=16)
tf.paragraphs[0].text = ""
for cat, tools in stacks:
    add_bullet(tf, f"{cat}:  {tools}", size=17)

# === SLIDE 7: Education & Certifications ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Education & Certifications", size=36, bold=True, color=DARK)
accent_bar(s)

tf = add_text(s, 0.8, 1.7, 5.5, 5, "", size=16)
tf.paragraphs[0].text = ""
add_bullet(tf, "Education", size=22, bold=True, color=TERRACOTTA)
add_bullet(tf, "Penn State University (2011–2015)", size=17, bold=True)
add_bullet(tf, "  BS, Real Estate Economics", size=15, color=GRAY)
add_bullet(tf, "UC Irvine (2018–Present)", size=17, bold=True)
add_bullet(tf, "  Data Analytics Bootcamp — Python, Tableau, BI Tools", size=15, color=GRAY)

tf2 = add_text(s, 6.8, 1.7, 5.5, 5, "", size=16)
tf2.paragraphs[0].text = ""
add_bullet(tf2, "Certifications", size=22, bold=True, color=TERRACOTTA)
add_bullet(tf2, "Intro to Enterprise-ready Gen AI Apps", size=17, bold=True)
add_bullet(tf2, "  Weaviate — Sep 2024", size=15, color=GRAY)
add_bullet(tf2, "", size=10)
add_bullet(tf2, "Courses", size=22, bold=True, color=TERRACOTTA)
add_bullet(tf2, "Cash Flow Modeling: MBS", size=17, bold=True)
add_bullet(tf2, "  Independent Study", size=15, color=GRAY)

# === SLIDE 8: Testimonial ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Client Testimonial", size=36, bold=True, color=DARK)
accent_bar(s)

add_text(s, 1.2, 2.0, 10.5, 3,
    '"Greg was a game-changer for our artist lead pipeline. He used the Spotify API to automate what used to be a super tedious process, saving us tons of time and hassle. Now, we\'ve got a system that not only finds potential leads but organizes them in a way that\'s easy to act on. Greg\'s technical know-how and out-of-the-box thinking really shined throughout the project, and he made everything so straightforward for us. Thanks to him, we can focus on connecting with artists instead of getting bogged down in data collection."',
    size=18, color=DARK)
add_text(s, 1.2, 5.2, 10.5, 1,
    "— Michael Miller, Artist Manager & Founder, Joywave Management\n   ex-Salesforce/DoorDash  |  November 27, 2024",
    size=16, bold=True, color=TERRACOTTA)

# === SLIDE 9: Personal Intelligence — Who Greg Is ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Personal Intelligence — Who Greg Is Off The Clock", size=32, bold=True, color=DARK)
accent_bar(s)

tf = add_text(s, 0.8, 1.7, 6.0, 5.2, "", size=16)
tf.paragraphs[0].text = ""
add_bullet(tf, "Voice & Language", size=19, bold=True, color=TERRACOTTA)
add_bullet(tf, 'Bio: "Me gusta mi voz en español más que inglés"', size=15)
add_bullet(tf, 'Captions in Spanish: "Agradecido", "Buscame aqui",', size=15)
add_bullet(tf, '  "nos vemos a la proxima", "excelente papito"', size=15)
add_bullet(tf, 'Italian: "Torino sei bellissima 🇮🇹" — Turin trip, Jul 2024', size=15)
add_bullet(tf, "Posts naturally in 3 languages — not performative", size=15)
add_bullet(tf, "", size=8)
add_bullet(tf, "Emotional Tone", size=19, bold=True, color=TERRACOTTA)
add_bullet(tf, '"Could not feel more grateful for these people', size=15)
add_bullet(tf, '  and the memories made" — Lightning in a Bottle', size=15)
add_bullet(tf, '"Much love lately 💕" — simple, present, open', size=15)
add_bullet(tf, 'Friends call him: "Big sunshine guy", "Greggy"', size=15)
add_bullet(tf, "Warm and magnetic — draws people in naturally", size=15)

tf2 = add_text(s, 7.0, 1.7, 5.8, 5.2, "", size=16)
tf2.paragraphs[0].text = ""
add_bullet(tf2, "Music & Nightlife", size=19, bold=True, color=TERRACOTTA)
add_bullet(tf2, '"Just a couple of plant guys at the afters"', size=15)
add_bullet(tf2, "After-party culture, house/electronic music scene", size=15)
add_bullet(tf2, '"Love when homies capture 📷 #CandidInTheMix"', size=15)
add_bullet(tf2, "DJ-adjacent — values unposed, in-the-moment shots", size=15)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "Aesthetics & Style", size=19, bold=True, color=TERRACOTTA)
add_bullet(tf2, "Festival: fur coats, tie-dye, bandanas, wristbands", size=15)
add_bullet(tf2, "Formal: sharp black suit, composed mirror selfie", size=15)
add_bullet(tf2, 'Moody/cinematic: candlelit smoking shot —', size=15)
add_bullet(tf2, '  friends joked "Narcos theme song", "Pablo Escobar"', size=15)
add_bullet(tf2, "Lives the duality fully — never one-dimensional", size=15)

# Source footnote slide 9
add_text(s, 0.8, 6.8, 11.5, 0.5,
    "Sources: /p/CuStB3aMpHZ  /p/C9UER8gO5l7  /p/C0KNxAXPpwl  /p/CsPCnk-vo5g  /p/CrBc8g3rGmW  /p/Cezex8-pe9R",
    size=9, color=GRAY)

# === SLIDE 10: Personal Intelligence — His World ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Personal Intelligence — His World, Events & Crew", size=32, bold=True, color=DARK)
accent_bar(s)

tf = add_text(s, 0.8, 1.7, 6.0, 5.2, "", size=16)
tf.paragraphs[0].text = ""
add_bullet(tf, "Annual Sacred Events", size=19, bold=True, color=TERRACOTTA)
add_bullet(tf, "Lightning in a Bottle (LiB) — California arts festival", size=15)
add_bullet(tf, '  "you\'ve done it again... see you next year" — deep devotion', size=15)
add_bullet(tf, "  Goes with a tight recurring crew every year", size=15)
add_bullet(tf, "Día de los Muertos — attends annually, captions in Spanish", size=15)
add_bullet(tf, "  'nos vemos a la próxima' — 'see you next time'", size=15)
add_bullet(tf, "Coachella-style festival (May 2022) — 'Buscame aqui'", size=15)
add_bullet(tf, "", size=8)
add_bullet(tf, "International Life", size=19, bold=True, color=TERRACOTTA)
add_bullet(tf, "Turin, Italy (Jul 2024) — Italian friends, 'see you next year'", size=15)
add_bullet(tf, "Santiago, Chile (2019) — lived there for work", size=15)
add_bullet(tf, "Valparaíso, Chile — Instagram highlight 'Valpo'", size=15)
add_bullet(tf, "Venice, LA — home base, street life, palm trees", size=15)

tf2 = add_text(s, 7.0, 1.7, 5.8, 5.2, "", size=16)
tf2.paragraphs[0].text = ""
add_bullet(tf2, "Close Circle (Recurring Names)", size=19, bold=True, color=TERRACOTTA)
add_bullet(tf2, "evanagony — close friend, calls him 'Al God'", size=15)
add_bullet(tf2, "natalie__harper — appears across multiple posts", size=15)
add_bullet(tf2, "collinjblack — 'Neighborhood swag' (possible family?)", size=15)
add_bullet(tf2, "lauraegalli — recurring across years of posts", size=15)
add_bullet(tf2, "jp.carroll — in his following, festival crew", size=15)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "Conversation Openers From Posts", size=19, bold=True, color=TERRACOTTA)
add_bullet(tf2, '"You go to Lightning in a Bottle every year?"', size=15)
add_bullet(tf2, '"What was Turin for — festival or just travel?"', size=15)
add_bullet(tf2, '"Your Día de los Muertos post was wild — where was that?"', size=15)
add_bullet(tf2, '"You DJ or just live in that world?"', size=15)

# Source footnote slide 10
add_text(s, 0.8, 6.8, 11.5, 0.5,
    "Sources: /p/CeWJQgELQc9 (LiB)  /p/Czt89wEPpo6 (Dia de Muertos)  /p/CdWE8JcJJ05 (Coachella)  /p/C9UER8gO5l7 (Turin)  instagram.com/gregorablack",
    size=9, color=GRAY)

# === SLIDE 11: Organizations & Interests ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Organizations & Influences", size=36, bold=True, color=DARK)
accent_bar(s)

tf = add_text(s, 0.8, 1.7, 5.5, 5, "", size=16)
tf.paragraphs[0].text = ""
add_bullet(tf, "Organizations", size=22, bold=True, color=TERRACOTTA)
add_bullet(tf, "Toastmasters International", size=17, bold=True)
add_bullet(tf, "  Nov 2015 – Present (10+ years)", size=15, color=GRAY)
add_bullet(tf, "Penn State Alumni Association", size=17, bold=True)
add_bullet(tf, "  Membership Committee Post-Grad Liaison", size=15, color=GRAY)
add_bullet(tf, "  Jan 2016 – Present", size=15, color=GRAY)

tf2 = add_text(s, 6.8, 1.7, 5.5, 5, "", size=16)
tf2.paragraphs[0].text = ""
add_bullet(tf2, "Top Voices Followed", size=22, bold=True, color=TERRACOTTA)
add_bullet(tf2, "Gary Vaynerchuk (2nd connection)", size=17, bold=True)
add_bullet(tf2, "  VaynerX, VaynerMedia, VeeFriends", size=15, color=GRAY)
add_bullet(tf2, "Kevin O'Leary (3rd connection)", size=17, bold=True)
add_bullet(tf2, "  O'Leary Ventures, Beanstox, Shark Tank", size=15, color=GRAY)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "Highlights", size=22, bold=True, color=TERRACOTTA)
add_bullet(tf2, "IG: Studio, 2023, Vibing, Valpo, Dalliance 7/18", size=15)

# === SLIDE 11b: DOGE Intelligence ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "⚠️  Critical Intel: DOGE Impact on AFS", size=34, bold=True, color=TERRACOTTA)
accent_bar(s)

tf = add_text(s, 0.8, 1.7, 11.5, 5.2, "", size=16)
tf.paragraphs[0].text = ""
add_bullet(tf, "What Happened", size=20, bold=True, color=DARK)
add_bullet(tf, "GSA directed agencies to review contracts with top 10 consulting firms", size=17)
add_bullet(tf, "30+ AFS contracts terminated  →  ~$240M in claimed DOGE savings", size=17, bold=True)
add_bullet(tf, "10 DOE task orders canceled alone  (~$92.8M)", size=17)
add_bullet(tf, "Pentagon canceled contracts affecting Accenture, Booz Allen, and Deloitte", size=17)
add_bullet(tf, "", size=8)
add_bullet(tf, "Business Impact", size=20, bold=True, color=DARK)
add_bullet(tf, "Accenture shares fell 7%+ after revenue warning", size=17)
add_bullet(tf, "New bookings down 3% to $20.9B in Q2", size=17)
add_bullet(tf, "Company maintains federal work is 'mission-critical'", size=17)
add_bullet(tf, "", size=8)
add_bullet(tf, "Why This Matters for Gregory", size=20, bold=True, color=TERRACOTTA)
add_bullet(tf, "He joined AFS in January 2026 — right as DOGE cuts hit hardest", size=17, bold=True)
add_bullet(tf, "AI roles may be protected (mission-critical) or at elevated risk", size=17)
add_bullet(tf, "GB Automation may serve as his hedge against AFS instability", size=17)
add_bullet(tf, "Great conversation angle: 'How is the DOGE situation affecting your work at AFS?'", size=17, color=GRAY)

# === SLIDE 12: Key Takeaways 1/2 ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Key Takeaways (1/2)", size=36, bold=True, color=DARK)
accent_bar(s)

takeaways1 = [
    "Real estate → data science → AI engineering → agentic systems entrepreneur. Non-linear, self-made path.",
    "Dual identity: Corporate stability (AFS Lead AI Engineer) + entrepreneurial hustle (GB Automation).",
    "Proven client results: Spotify API automation for Joywave, 40% conversion AI SDR, UK FinTech sales ops.",
    "Full-stack AI consulting: Claude API, LangChain, RAG, n8n, Supabase, Salesforce, FastAPI — rare breadth.",
    "International experience: Lived/worked in Santiago, Chile. Bilingual English/Spanish.",
]
tf = add_text(s, 0.8, 1.7, 11.5, 5, "", size=16)
tf.paragraphs[0].text = ""
for i, t in enumerate(takeaways1, 1):
    add_bullet(tf, f"{i}. {t}", size=17)

# === SLIDE 13: Key Takeaways 2/2 ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Key Takeaways (2/2)", size=36, bold=True, color=DARK)
accent_bar(s)

takeaways2 = [
    "~3M Threads followers vs 1.4K IG — massive audience exists, just not on Instagram. Threads is the distribution engine.",
    "10+ years in Toastmasters — strong communicator and public speaker. Selling should come naturally.",
    "Real estate foundation (Penn State + 3 years CRE) gives unique domain knowledge for proptech/fintech clients.",
    "Follows Gary Vee and Kevin O'Leary — entrepreneurial, growth-minded, values hustle culture.",
    "Instagram is intentionally personal (no business content) — deliberate brand separation from professional identity.",
]
tf = add_text(s, 0.8, 1.7, 11.5, 5, "", size=16)
tf.paragraphs[0].text = ""
for i, t in enumerate(takeaways2, 6):
    add_bullet(tf, f"{i}. {t}", size=17)

# === SLIDE 14: Conversation Starters ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Conversation Starters", size=36, bold=True, color=DARK)
accent_bar(s)

starters = [
    "The Spotify API project for Joywave — how did that client relationship start?",
    "Transition from real estate to data science — what was the catalyst moment?",
    "Chile experience at Econsult — was that what sparked the Spanish fluency?",
    "Threads growth to ~3M — what content resonated? Was it organic?",
    "Running GB Automation alongside Accenture during the DOGE cuts — how are you navigating that?",
    "The 40% conversion rate on the AI SDR Campaign Manager — what made it work?",
    "gbautomation.xyz is under construction — what's the vision for the website launch?",
    "Toastmasters for 10+ years — how has that shaped your consulting sales approach?",
]
tf = add_text(s, 0.8, 1.7, 11.5, 5, "", size=16)
tf.paragraphs[0].text = ""
for i, t in enumerate(starters, 1):
    add_bullet(tf, f"{i}. {t}", size=18)

# === SLIDE 15: Next Steps ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, 0.8, 0.5, 11, 0.7, "Next Steps", size=36, bold=True, color=DARK)
accent_bar(s)

tf = add_text(s, 0.8, 1.7, 11.5, 5, "", size=16)
tf.paragraphs[0].text = ""
add_bullet(tf, "Profile Audit Actions", size=22, bold=True, color=TERRACOTTA)
add_bullet(tf, "Update LinkedIn location (Venice LA vs New York discrepancy)", size=17)
add_bullet(tf, "Add AI/automation content to Instagram (or keep separation intentional)", size=17)
add_bullet(tf, "Leverage Threads audience for consulting lead gen", size=17)
add_bullet(tf, "Request more LinkedIn recommendations (only 1 currently)", size=17)
add_bullet(tf, "Fill in Sr Data and AI Engineer role details (partially visible)", size=17)
add_bullet(tf, "", size=10)
add_bullet(tf, "Research Follow-ups", size=22, bold=True, color=TERRACOTTA)
add_bullet(tf, "Deep-dive @gb.noir.lugar secondary account", size=17)
add_bullet(tf, "Threads content analysis (what drove 3M followers?)", size=17)
add_bullet(tf, "AFS internal visibility — is there an AI practice lead opportunity?", size=17)

# Save
out = r"C:\Users\gblac\OneDrive\Desktop\consulting-co\.claude\context\clients\gregory-black\Gregory Black - Client Research v2.pptx"
prs.save(out)
print(f"Saved: {out}")
